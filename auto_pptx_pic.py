from matplotlib.pyplot import colormaps
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.util import Cm, Pt
import os
from posixpath import normcase


# Presentation オブジェクトを生成
ppt = Presentation()

# スライドのサイズを指定
ppt.slide_width = Cm(33.868)
ppt.slide_heght = Cm(19.05)

# ----------------------------------------------------------------------
# スライド1枚目
# 追加するスライドを選択
slide_layout_0 = ppt.slide_layouts[0]
# スライドを追加
slide_0 = ppt.slides.add_slide(slide_layout_0)
# place holder (テキスト入れるやつ？)を選択
slide_0_title = slide_0.placeholders[0]
# スライドに貼り付け
print(type(slide_0_title))
# スライドの幅
slide_0_title.width = Cm(30)
# スライドの高さ
slide_0_title.height = Cm(20)
# スライドの文字
slide_0_title.text = "pic"


status = ["NSE", "freezeout", "last"]

Ye = ["010", "013", "016"]

T0 = ["4e9", "7e9", \
    "1e10", "4e10", "7e10", \
        "1e11", "4e11", "7e11", \
            "1e12"]
rho0 = ["1e10", "4e10", "7e10", \
    "1e11", "4e11", "7e11", \
        "1e12", "4e12", "7e12", \
            "1e13", "4e13"]

for w in status:
    for x in Ye:
        n = 0
        for y in T0:
            m = 0
            for z in rho0:
                # ----------------------------------------------------------------------
                # スライド2枚目以降
                #
                picname = w + "_Ye_" + x + "_T0_" + y + "_rho0_" + z + ".png"
                input_pic_path = "../../progress/0924to1004/pic_edit/" + picname
                exsist = os.path.exists(input_pic_path)
                if not exsist:
                    m += 1
                    continue

                # スライドのレイアウトを選択
                slide_layout_6 = ppt.slide_layouts[6]
                # スライドを追加
                slide_6 = ppt.slides.add_slide(slide_layout_6)

                # 画像のplaceholder を作成
                shapes = slide_6.shapes
                shapes.add_picture(input_pic_path, Cm(15), Cm(3), width = None, height = Cm(10.88))
                colormap_path = "../../progress/0924to1004/Thorium232/last_Thorium232_Ye_" + x + ".png"
                shapes.add_picture(colormap_path, 0, Cm(3), width = None, height = Cm(10.88))

                # 指定する長方形を追加
                rect0 = slide_6.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    Pt(52 + 23 * m), Pt(122 + 26 * n),
                    Pt(22), Pt(26)
                )
                rect0.fill.background()
                rect0.line.width = Pt(2)

                # # placeholder_1 = slide_6.placeholders[1]
                shape = shapes.add_textbox(Cm(20), Cm(14), Cm(22), Cm(16))
                # # box内にテキストを追加
                shape.text = picname
                # ----------------------------------------------------------------------
                m += 1
            n += 1
            # スライドのレイアウトを選択
            slide_layout_6 = ppt.slide_layouts[6]
            # スライドを追加
            slide_6 = ppt.slides.add_slide(slide_layout_6)


        m = 0
        for z in rho0:
            n = 0
            for y in T0:
                # ----------------------------------------------------------------------
                # スライド2枚目以降
                #
                picname = w + "_Ye_" + x + "_T0_" + y + "_rho0_" + z + ".png"
                input_pic_path = "../../progress/0924to1004/pic_edit/" + picname
                exsist = os.path.exists(input_pic_path)
                if not exsist:
                    m += 1
                    continue

                # スライドのレイアウトを選択
                slide_layout_6 = ppt.slide_layouts[6]
                # スライドを追加
                slide_6 = ppt.slides.add_slide(slide_layout_6)

                # 画像のplaceholder を作成
                shapes = slide_6.shapes
                shapes.add_picture(input_pic_path, Cm(15), Cm(3), width = None, height = Cm(10.88))
                colormap_path = "../../progress/0924to1004/Thorium232/last_Thorium232_Ye_" + x + ".png"
                shapes.add_picture(colormap_path, 0, Cm(3), width = None, height = Cm(10.88))

                # 指定する長方形を追加
                rect0 = slide_6.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    Pt(52 + 23 * m), Pt(122 + 26 * n),
                    Pt(22), Pt(26)
                )
                rect0.fill.background()
                rect0.line.width = Pt(2)

                # # placeholder_1 = slide_6.placeholders[1]
                shape = shapes.add_textbox(Cm(20), Cm(14), Cm(22), Cm(16))
                # # box内にテキストを追加
                shape.text = picname
                # ----------------------------------------------------------------------
                n += 1
            m += 1
            # スライドのレイアウトを選択
            slide_layout_6 = ppt.slide_layouts[6]
            # スライドを追加
            slide_6 = ppt.slides.add_slide(slide_layout_6)
        # スライドのレイアウトを選択
        slide_layout_6 = ppt.slide_layouts[6]
        # スライドを追加
        slide_6 = ppt.slides.add_slide(slide_layout_6)
    # スライドのレイアウトを選択
    slide_layout_6 = ppt.slide_layouts[6]
    # スライドを追加
    slide_6 = ppt.slides.add_slide(slide_layout_6)



# 保存
ppt.save("../../progress/0924to1004/analysis/pic.pptx")